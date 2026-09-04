"""tests/test_slides.py — slides capability contracts (P3).

Pins: PPTD parse/validate (the zero-LLM gate the author loop repairs
against), theme resolution, the rich-text parser, the HTML renderer's
determinism, the PPTX exporter's structural guarantees (real shapes, fade
order, CRC), the recipe's quality discipline (bad pages retain retry
diagnostics but block publication), and the runtime/starter kind contract that keeps /api/v1/tasks/*
able to poll what /start launches.
"""

from __future__ import annotations

import os
import sys

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import lib.slides.pptd as pptd  # noqa: E402
from lib.slides.pptd import Deck, Page, parse_deck, validate_deck  # noqa: E402

pytestmark = pytest.mark.unit


# ── Fixtures ──────────────────────────────────────────────

def _write_deck(tmp_path, pages_yaml: list, *, theme=None, size=(1280, 720)):
    """Write a minimal deck dir; return manifest path."""
    import yaml
    deck_dir = tmp_path / 'deck'
    (deck_dir / 'pages').mkdir(parents=True)
    names = []
    for i, text in enumerate(pages_yaml, 1):
        name = f'pages/{i:02d}.page'
        (deck_dir / name).write_text(text, encoding='utf-8')
        names.append(name)
    manifest = {'version': 'v2', 'title': '测试 deck', 'size': list(size),
                'theme': theme or {
                    'colors': {'bg': '#F7F7F5', 'ink': '#1B2430',
                               'primary': '#16283C', 'accent': '#C0652B',
                               'muted': '#6B7280', 'hairline': '#D8D5CE'},
                    'textStyles': {'title': {'fontSize': 40,
                                             'color': '$primary'},
                                   'body': {'fontSize': 18, 'color': '$ink'}},
                },
                'pages': names}
    (deck_dir / 'deck.pptd').write_text(
        yaml.safe_dump(manifest, allow_unicode=True), encoding='utf-8')
    return str(deck_dir / 'deck.pptd')


_COVER = '''pageType: cover
background: {type: solid, color: "$bg"}
elements:
  - elementId: title
    elementType: text
    bounds: [72, 200, 1136, 120]
    content:
      style: "$title"
      align: [left, middle]
      text: |
        <p><strong>一寸万象</strong></p>
  - elementId: rule
    elementType: shape
    bounds: [72, 340, 64, 6]
    shapeName: rect
    fill: {type: solid, color: "$accent"}
'''

_TABLE_PAGE = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: t
    elementType: table
    bounds: [80, 120, 1120, 300]
    columnWidths: [0.5, 0.5]
    rowHeights: [0.34, 0.33, 0.33]
    rows:
      - - text: "指标"
        - text: "2025"
      - - text: "营收"
        - text: "96.3"
      - - text: "利润"
        - text: "15.8"
'''


# ── parse / validate ──────────────────────────────────────

class TestParseValidate:
    def test_round_trip_clean(self, tmp_path):
        deck = parse_deck(_write_deck(tmp_path, [_COVER, _TABLE_PAGE]))
        assert deck.title == '测试 deck'
        assert deck.size == (1280, 720)
        assert len(deck.pages) == 2
        assert validate_deck(deck) == []

    def test_manifest_dir_resolution(self, tmp_path):
        path = _write_deck(tmp_path, [_COVER])
        deck = parse_deck(os.path.dirname(path))
        assert len(deck.pages) == 1

    def test_bad_version_rejected(self, tmp_path):
        path = _write_deck(tmp_path, [_COVER])
        text = open(path).read().replace('v2', 'v9')
        open(path, 'w').write(text)
        with pytest.raises(pptd.PPTDError):
            parse_deck(path)

    def test_path_escape_rejected(self, tmp_path):
        deck_dir = tmp_path / 'deck'
        (deck_dir / 'pages').mkdir(parents=True)
        (deck_dir / 'deck.pptd').write_text(
            'version: v2\nsize: [1280, 720]\npages: ["../escape.page"]\n',
            encoding='utf-8')
        with pytest.raises(pptd.PPTDError):
            parse_deck(str(deck_dir / 'deck.pptd'))

    def test_validator_catches_real_defects(self, tmp_path):
        bad = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: a
    elementType: text
    bounds: [10, 10, 100, 50]
    content: {style: "$nosuch", text: "x"}
  - elementId: a
    elementType: text
    bounds: [10, 10, 100, 50]
    content: {text: "dup id"}
  - elementId: b
    elementType: shape
    bounds: [10, 10, 100, 50]
    shapeName: nonExistentShape
  - elementId: c
    elementType: table
    bounds: [10, 10, 100, 50]
    columnWidths: [0.5, 0.6]
    rowHeights: [1.0]
    rows: [[{text: "x"}, {text: "y"}]]
'''
        deck = parse_deck(_write_deck(tmp_path, [bad]))
        findings = validate_deck(deck)
        blob = '\n'.join(findings)
        assert 'unknown textStyle token' in blob
        assert 'duplicate elementId' in blob
        assert 'unsupported shape' in blob
        assert 'columnWidths must sum to 1' in blob

    def test_allow_overlap_must_be_explicit_boolean(self, tmp_path):
        page = _COVER.replace('elementType: text',
                              'elementType: text\n    allowOverlap: yes-please', 1)
        findings = validate_deck(parse_deck(_write_deck(tmp_path, [page])))
        assert any('allowOverlap must be boolean' in f for f in findings)

    def test_text_fit_mode_is_validated(self, tmp_path):
        page = _COVER.replace('style: "$title"',
                              'style: "$title"\n      fit: overflow', 1)
        findings = validate_deck(parse_deck(_write_deck(tmp_path, [page])))
        assert any('text fit must be shrink|none|resize' in f
                   for f in findings)

    def test_theme_resolution(self, tmp_path):
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        assert pptd.resolve_color('$primary', deck.theme) == '#16283C'
        assert pptd.resolve_color('#AABBCCDD', deck.theme) == '#AABBCCDD'
        assert pptd.resolve_color('$nope', deck.theme, 'x') == 'x'
        st = pptd.text_style({'style': '$title'}, deck.theme)
        assert st['fontSize'] == 40 and st['color'] == '#16283C'
        st2 = pptd.text_style({'style': '$title', 'fontSize': 22},
                              deck.theme)
        assert st2['fontSize'] == 22      # inline wins over the theme style


# ── rich text ─────────────────────────────────────────────

class TestRichText:
    def test_plain_shorthand(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text('第一行\n第二行', {})
        assert len(paras) == 2
        assert paras[0].runs[0].text == '第一行'

    def test_marks_and_inline_styles(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text(
            '<p style="text-align:center"><strong>重</strong>'
            '<span style="color:$primary;font-size:24px">点</span></p>',
            {'colors': {'primary': '#123456'}})
        runs = paras[0].runs
        assert paras[0].align == 'center'
        assert runs[0].bold and runs[0].text == '重'
        assert runs[1].color == '#123456' and runs[1].font_size == 24

    def test_lists(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text('<ul><li>甲</li><li>乙</li></ul>', {})
        assert len(paras) == 2 and paras[0].list_kind == 'ul'

    def test_malformed_degrades_to_text(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text('<p><strong>未闭合', {})
        assert paras and '未闭合' in paras[0].runs[0].text

    def test_pretty_printed_sibling_paragraphs_do_not_create_blank_lines(self):
        from lib.slides.richtext import parse_rich_text
        paras = parse_rich_text(
            '<p>澎程不是子品牌,</p>\n'
            '        <p>是小米汽车旗下平级的</p>\n'
            '        <p>第二大产品系列。</p>', {})
        assert len(paras) == 3
        assert [''.join(r.text for r in p.runs) for p in paras] == [
            '澎程不是子品牌,', '是小米汽车旗下平级的', '第二大产品系列。']


# ── HTML renderer ─────────────────────────────────────────

class TestRenderHtml:
    def test_page_html_deterministic_and_themed(self, tmp_path):
        from lib.slides.render_html import render_page_html
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        a = render_page_html(deck, deck.pages[0])
        b = render_page_html(deck, deck.pages[0])
        assert a == b                       # deterministic
        assert '#F7F7F5' in a               # $bg resolved
        assert '#16283C' in a               # $primary resolved
        assert '1280px' in a and '720px' in a
        assert '一寸万象' in a
        assert 'data-element-id="title"' in a
        assert 'data-allow-overlap="false"' in a

    def test_gradient_angle_mapping(self):
        from lib.slides.render_html import _gradient_css
        css = _gradient_css({'type': 'gradient', 'gradientType': 'linear',
                             'angle': 90,
                             'stops': [{'position': 0, 'color': '#000000'},
                                       {'position': 1, 'color': '#FFFFFF'}]},
                            {})
        assert '180.0deg' in css            # PPTD 90 (top→bottom) = CSS 180

    def test_table_and_richtext_render(self, tmp_path):
        from lib.slides.render_html import render_page_html
        deck = parse_deck(_write_deck(tmp_path, [_TABLE_PAGE]))
        html = render_page_html(deck, deck.pages[0])
        assert '<table' in html and '指标' in html and 'colspan' not in html


# ── PPTX exporter ─────────────────────────────────────────

class TestExportPptx:
    def test_export_structure(self, tmp_path):
        pytest.importorskip('pptx')
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER, _TABLE_PAGE]))
        out = str(tmp_path / 'out.pptx')
        summary = export_pptx(deck, out)
        assert summary['slides'] == 2
        assert summary['fadeTransitions'] == 2
        assert summary['bytes'] > 4096

        from pptx import Presentation
        prs = Presentation(out)
        assert len(prs.slides) == 2
        assert prs.slide_width == 1280 * 12700
        s1, s2 = prs.slides
        texts = []
        for sh in s1.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
        assert any('一寸万象' in t for t in texts)
        assert any(sh.has_table for sh in s2.shapes)

    def test_fade_transition_order(self, tmp_path):
        pytest.importorskip('pptx')
        import re
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        with zipfile.ZipFile(out) as z:
            xml = z.read('ppt/slides/slide1.xml').decode()
        i_csld = xml.index('</p:cSld>')
        i_trans = xml.index('<p:transition')
        assert i_trans > i_csld            # CT_Slide order
        assert xml.count('<p:fade/>') == 1

    def test_notes_and_font_ea(self, tmp_path):
        pytest.importorskip('pptx')
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        with zipfile.ZipFile(out) as z:
            xml = z.read('ppt/slides/slide1.xml').decode()
        assert '<a:ea typeface=' in xml     # CJK font on the ea slot

    def test_text_boxes_are_fixed_shrink_to_fit_with_zero_para_spacing(
            self, tmp_path):
        pytest.importorskip('pptx')
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        with zipfile.ZipFile(out) as z:
            xml = z.read('ppt/slides/slide1.xml').decode()
        assert '<a:normAutofit' in xml
        assert '<a:spAutoFit' not in xml
        assert '<a:spcBef><a:spcPts val="0"/></a:spcBef>' in xml
        assert '<a:spcAft><a:spcPts val="0"/></a:spcAft>' in xml


class TestLayoutQA:
    def test_real_line_rect_collision_overflow_and_escape_hatch(self):
        from lib.slides.layout_qa import _page_findings
        records = [
            {'id': 'overflowing', 'visible': True, 'allowOverlap': False,
             'outer': {'left': 0, 'top': 0, 'right': 100, 'bottom': 30},
             'rects': [{'left': 2, 'top': 2, 'right': 120, 'bottom': 24}]},
            {'id': 'collision', 'visible': True, 'allowOverlap': False,
             'outer': {'left': 80, 'top': 0, 'right': 180, 'bottom': 30},
             'rects': [{'left': 82, 'top': 2, 'right': 150, 'bottom': 24}]},
            {'id': 'decorative', 'visible': True, 'allowOverlap': True,
             'outer': {'left': 0, 'top': 0, 'right': 180, 'bottom': 30},
             'rects': [{'left': 0, 'top': 2, 'right': 180, 'bottom': 24}]},
        ]
        findings = _page_findings(records, 0, tolerance=1.5)
        kinds = [f['type'] for f in findings]
        assert kinds.count('text_overflow') == 1
        assert kinds.count('text_collision') == 1
        assert findings[-1]['elements'] == ['overflowing', 'collision']

    def test_powerpoint_metric_reserve_and_later_image_occlusion(self):
        from lib.slides.layout_qa import _page_findings
        title = {
            'id': 'title', 'z': 0, 'visible': True, 'allowOverlap': False,
            'fontSize': 40,
            'outer': {'left': 20, 'top': 20, 'right': 500, 'bottom': 120},
            'rects': [
                {'left': 20, 'top': 24, 'right': 300, 'bottom': 64},
                {'left': 20, 'top': 68, 'right': 340, 'bottom': 112},
            ],
        }
        image = {
            'id': 'hero', 'type': 'image', 'z': 1, 'visible': True,
            'outer': {'left': 20, 'top': 124, 'right': 500, 'bottom': 400},
        }
        findings = _page_findings([title], 0, tolerance=1.5,
                                  occluders=[image])
        assert [f['type'] for f in findings] == [
            'pptx_text_overflow_risk', 'text_image_occlusion_risk']
        assert findings[-1]['elements'] == ['title', 'hero']


# ── Recipe (mocked LLM) ───────────────────────────────────

class TestRecipe:
    def _fake_llm(self, content):
        def _call(messages, **kw):
            return content, {'total_tokens': 10}
        return _call

    def test_page_author_discards_late_reply_and_propagates_policy(
            self, tmp_path, monkeypatch):
        import threading

        import lib.slides.author as author

        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        abort_event = threading.Event()
        seen = {}

        def late_reply(messages, **kwargs):
            seen.update(kwargs)
            abort_event.set()
            return 'late yaml', {}

        monkeypatch.setattr(author, '_llm', late_reply)
        result = author.author_page(
            deck, {'pageType': 'cover', 'key_message': '标题'}, 0, 1,
            abort_check=abort_event.is_set, max_429_attempts=7)

        assert result['mode'] == 'aborted'
        assert result['yaml'] == ''
        assert seen['abort_check']() is True
        assert seen['max_429_attempts'] == 7

    @staticmethod
    def _author_context(tmp_path, page_count=4, *, abort_event=None,
                        emit=None):
        manifest = _write_deck(tmp_path, [_COVER])
        seed_deck = parse_deck(manifest)
        page_types = ['cover'] + ['content'] * max(0, page_count - 2)
        if page_count > 1:
            page_types.append('final')
        pages = [
            {'pageType': page_type, 'purpose': f'purpose {index}',
             'key_message': f'message {index}', 'layout_hint': 'editorial',
             'content_notes': f'notes {index}'}
            for index, page_type in enumerate(page_types)
        ]
        return {
            'deck_dir': seed_deck.root,
            'size': seed_deck.size,
            'lang': 'zh',
            'model': 'test-model',
            'author_rounds': 1,
            'abort_event': abort_event,
            'emit': emit,
            'max_429_attempts': 7,
            'artifacts': {
                'outline': {'title': 'Author cache deck', 'pages': pages},
                'design': {'theme_id': 'paper-engineer',
                           'theme_tokens': seed_deck.theme},
            },
        }

    @staticmethod
    def _authored_page_yaml(label):
        return (
            'pageType: content\n'
            'background: {type: solid, color: "$bg"}\n'
            'elements:\n'
            '  - elementId: title\n'
            '    elementType: text\n'
            '    bounds: [72, 72, 1136, 120]\n'
            '    content:\n'
            '      style: "$title"\n'
            '      text: |\n'
            f'        <p>{label}</p>\n')

    def test_author_pages_use_bounded_fanout_and_stable_order(
            self, tmp_path, monkeypatch):
        import threading

        import lib.slides.author as author
        import lib.slides.recipe as recipe

        real_prepare = author.prepare_author_prompt_context
        prepared = []

        def prepare_once(*args, **kwargs):
            context = real_prepare(*args, **kwargs)
            prepared.append(context)
            return context

        monkeypatch.setattr(
            author, 'prepare_author_prompt_context', prepare_once)
        monkeypatch.setenv('TOFU_PRODUCTION_LLM_FANOUT', '2')
        monkeypatch.setattr(
            'lib.slides._asset_preflight.prepare_deck_assets',
            lambda *a, **k: {'by_page': {}, 'records': [], 'findings': []})
        lock = threading.Lock()
        release = threading.Event()
        state = {'active': 0, 'peak': 0, 'calls': []}
        prompt_context_ids = []

        def fake_author(_deck, _brief, index, _total, **_kwargs):
            with lock:
                state['active'] += 1
                state['peak'] = max(state['peak'], state['active'])
                state['calls'].append(index)
                prompt_context_ids.append(id(_kwargs['prompt_context']))
                if state['active'] == 2:
                    release.set()
            try:
                assert release.wait(3)
                return {'ok': True,
                        'yaml': self._authored_page_yaml(f'page-{index}'),
                        'mode': 'authored', 'rounds': 1, 'findings': []}
            finally:
                with lock:
                    state['active'] -= 1

        monkeypatch.setattr(author, 'author_page', fake_author)
        events = []
        ctx = self._author_context(tmp_path, emit=events.append)

        result = recipe._run_author(ctx)

        assert state['peak'] == 2
        assert len(prepared) == 1
        assert set(prompt_context_ids) == {id(prepared[0])}
        assert sorted(state['calls']) == [0, 1, 2, 3]
        assert result['page_files'] == [
            'pages/01_cover.page', 'pages/02_content.page',
            'pages/03_content.page', 'pages/04_final.page']
        for index, name in enumerate(result['page_files']):
            text = (tmp_path / 'deck' / name).read_text(encoding='utf-8')
            assert f'page-{index}' in text
        assert sorted(event['page'] for event in events) == [1, 2, 3, 4]

    def test_author_cache_reuses_exact_pages_and_repairs_only_misses(
            self, tmp_path, monkeypatch):
        import lib.slides.author as author
        import lib.slides.recipe as recipe

        monkeypatch.setenv('TOFU_PRODUCTION_LLM_FANOUT', '1')
        monkeypatch.setattr(
            'lib.slides._asset_preflight.prepare_deck_assets',
            lambda *a, **k: {'by_page': {}, 'records': [], 'findings': []})
        calls = []

        def fake_author(_deck, brief, index, _total, **_kwargs):
            calls.append(index)
            return {'ok': True,
                    'yaml': self._authored_page_yaml(
                        f'{index}-{brief["key_message"]}'),
                    'mode': 'authored', 'rounds': 1, 'findings': []}

        monkeypatch.setattr(author, 'author_page', fake_author)
        ctx = self._author_context(tmp_path, page_count=3)
        first = recipe._run_author(ctx)
        assert calls == [0, 1, 2]

        calls.clear()
        second = recipe._run_author(ctx)
        assert calls == []
        assert second['page_files'] == first['page_files']
        assert second['authored'] == 3

        ctx['artifacts']['outline']['pages'][1]['key_message'] = 'changed'
        calls.clear()
        recipe._run_author(ctx)
        assert calls == [1]

        page = tmp_path / 'deck' / first['page_files'][2]
        damaged = bytearray(page.read_bytes())
        damaged[-2] = ord('X') if damaged[-2] != ord('X') else ord('Y')
        page.write_bytes(damaged)
        calls.clear()
        recipe._run_author(ctx)
        assert calls == [2]

    def test_author_abort_stops_new_page_admission(
            self, tmp_path, monkeypatch):
        import threading

        import lib.slides.author as author
        import lib.slides.recipe as recipe
        from lib.production.stages import StageAborted

        monkeypatch.setenv('TOFU_PRODUCTION_LLM_FANOUT', '1')
        monkeypatch.setattr(
            'lib.slides._asset_preflight.prepare_deck_assets',
            lambda *a, **k: {'by_page': {}, 'records': [], 'findings': []})
        abort_event = threading.Event()
        calls = []

        def abort_first(_deck, _brief, index, _total, **_kwargs):
            calls.append(index)
            abort_event.set()
            return {'ok': False, 'yaml': '', 'mode': 'aborted', 'rounds': 0,
                    'findings': ['aborted']}

        monkeypatch.setattr(author, 'author_page', abort_first)
        ctx = self._author_context(tmp_path, abort_event=abort_event)

        with pytest.raises(StageAborted):
            recipe._run_author(ctx)
        assert calls == [0]
        assert not (tmp_path / 'deck' / 'pages/01_cover.page').exists()

    def test_author_fallback_is_not_reused(self, tmp_path, monkeypatch):
        import lib.slides.author as author
        import lib.slides.recipe as recipe

        monkeypatch.setenv('TOFU_PRODUCTION_LLM_FANOUT', '1')
        monkeypatch.setattr(
            'lib.slides._asset_preflight.prepare_deck_assets',
            lambda *a, **k: {'by_page': {}, 'records': [], 'findings': []})
        calls = []

        def fake_author(_deck, _brief, index, _total, **_kwargs):
            calls.append(index)
            return {'ok': True, 'yaml': self._authored_page_yaml(index),
                    'mode': 'fallback', 'rounds': 1, 'findings': ['test']}

        monkeypatch.setattr(author, 'author_page', fake_author)
        ctx = self._author_context(tmp_path, page_count=3)
        recipe._run_author(ctx)
        recipe._run_author(ctx)

        assert calls == [0, 1, 2, 0, 1, 2]

    @staticmethod
    def _visual_qa_context(tmp_path, *, page_count=4, abort_event=None):
        manifest = _write_deck(tmp_path, [_COVER] * page_count)
        deck = parse_deck(manifest)
        preview_dir = tmp_path / 'deck' / 'preview' / 'pages'
        preview_dir.mkdir(parents=True)
        previews = []
        for index in range(page_count):
            path = preview_dir / f'{index + 1:02d}.png'
            path.write_bytes(b'preview')
            previews.append(str(path))
        return {
            'deck_dir': deck.root,
            'lang': 'zh',
            'model': 'test-model',
            'qa_model': 'test-vlm',
            'abort_event': abort_event,
            'max_429_attempts': 7,
            'artifacts': {
                'outline': {'pages': [
                    {'pageType': 'content', 'key_message': f'page {index}'}
                    for index in range(page_count)
                ]},
                'design': {'theme_id': 'paper-engineer'},
                'author': {'assets_by_page': {}, 'input_images': []},
                'render': {'previews': previews},
            },
        }

    def test_visual_page_reviews_use_bounded_fanout(self, tmp_path,
                                                    monkeypatch):
        import threading

        import lib.design_sys.visual_qa as vqa
        import lib.slides.recipe as recipe

        monkeypatch.setenv('TOFU_PRODUCTION_LLM_FANOUT', '2')
        monkeypatch.setattr(vqa, 'visual_qa_available', lambda: (True, ''))
        monkeypatch.setattr(
            'lib.design_sys.contact_sheet.build_contact_sheet',
            lambda _previews, output, **_kwargs: output)
        lock = threading.Lock()
        release = threading.Event()
        state = {'active': 0, 'peak': 0, 'pages': []}

        def qa_frame(_path, **kwargs):
            label = kwargs['label']
            if label == 'deck-coherence':
                return {'ok': True, 'findings': []}
            with lock:
                state['active'] += 1
                state['peak'] = max(state['peak'], state['active'])
                state['pages'].append(label)
                if state['active'] == 2:
                    release.set()
            try:
                assert release.wait(3)
                return {'ok': True, 'findings': []}
            finally:
                with lock:
                    state['active'] -= 1

        monkeypatch.setattr(vqa, 'qa_frame', qa_frame)
        result = recipe._run_visual_qa(self._visual_qa_context(tmp_path))

        assert state['peak'] == 2
        assert sorted(state['pages']) == [
            'page-01', 'page-02', 'page-03', 'page-04']
        assert result['clean'] == 4
        assert result['repaired'] == 0

    def test_visual_page_abort_stops_new_review_admission(
            self, tmp_path, monkeypatch):
        import threading

        import lib.design_sys.visual_qa as vqa
        import lib.slides.recipe as recipe
        from lib.production.stages import StageAborted

        abort_event = threading.Event()
        monkeypatch.setenv('TOFU_PRODUCTION_LLM_FANOUT', '1')
        monkeypatch.setattr(vqa, 'visual_qa_available', lambda: (True, ''))
        monkeypatch.setattr(
            'lib.design_sys.contact_sheet.build_contact_sheet',
            lambda _previews, output, **_kwargs: output)
        page_calls = []

        def qa_frame(_path, **kwargs):
            if kwargs['label'] == 'deck-coherence':
                return {'ok': True, 'findings': []}
            page_calls.append(kwargs['label'])
            abort_event.set()
            return {'ok': False, 'skipped': True,
                    'reason': 'aborted after visual QA', 'findings': []}

        monkeypatch.setattr(vqa, 'qa_frame', qa_frame)
        ctx = self._visual_qa_context(
            tmp_path, abort_event=abort_event)

        with pytest.raises(StageAborted):
            recipe._run_visual_qa(ctx)
        assert page_calls == ['page-01']

    def test_visual_qa_cache_reuses_exact_pixels_and_reruns_only_changes(
            self, tmp_path, monkeypatch):
        from pathlib import Path

        import lib.design_sys.visual_qa as vqa
        import lib.slides.recipe as recipe

        monkeypatch.setenv('TOFU_PRODUCTION_LLM_FANOUT', '2')
        monkeypatch.setattr(vqa, 'visual_qa_available', lambda: (True, ''))

        def contact_sheet(previews, output, **_kwargs):
            path = Path(output)
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(b'|'.join(Path(item).read_bytes()
                                       for item in previews))
            return output

        monkeypatch.setattr(
            'lib.design_sys.contact_sheet.build_contact_sheet', contact_sheet)
        calls = []

        def qa_frame(path, **kwargs):
            calls.append(kwargs['label'])
            identity = vqa.qa_frame_input_sha256(
                path, theme=kwargs.get('theme'),
                subject=kwargs.get('subject', '视频帧'),
                model=kwargs.get('model', ''),
                max_tokens=kwargs.get('max_tokens', 1500))
            return {'ok': True, 'skipped': False, 'reason': '',
                    'findings': [], 'has_blocker': False,
                    'summary': '0 finding(s)',
                    'input_sha256': identity}

        monkeypatch.setattr(vqa, 'qa_frame', qa_frame)
        ctx = self._visual_qa_context(tmp_path)
        first = recipe._run_visual_qa(ctx)
        assert first['clean'] == 4
        assert set(calls) == {
            'deck-coherence', 'page-01', 'page-02', 'page-03', 'page-04'}

        calls.clear()
        second = recipe._run_visual_qa(ctx)
        assert second['clean'] == 4
        assert calls == []

        Path(ctx['artifacts']['render']['previews'][2]).write_bytes(
            b'changed-preview')
        calls.clear()
        third = recipe._run_visual_qa(ctx)
        assert third['clean'] == 4
        assert set(calls) == {'deck-coherence', 'page-03'}

    def test_full_graph_and_resume(self, tmp_path, monkeypatch):
        import lib.slides.recipe as recipe
        outline_json = ('{"title": "T", "scenario": "tech-engineering",'
                        ' "pages": ['
                        ' {"pageType": "cover", "purpose": "p", '
                        '  "key_message": "标题判断", "layout_hint": "", '
                        '  "content_notes": "n"},'
                        ' {"pageType": "content", "purpose": "p2", '
                        '  "key_message": "第二页判断", "layout_hint": "", '
                        '  "content_notes": "n2"},'
                        ' {"pageType": "final", "purpose": "p3", '
                        '  "key_message": "结论", "layout_hint": "", '
                        '  "content_notes": "n3"}]}')
        calls = {'n': 0}

        def _llm(messages, **kw):
            calls['n'] += 1
            text = messages[0]['content']
            if '大纲' in text or 'outline' in text:
                return outline_json, {}
            # page author path: return a minimal valid page
            return ('pageType: content\n'
                    'background: {type: solid, color: "$bg"}\n'
                    'elements:\n'
                    '  - elementId: t\n'
                    '    elementType: text\n'
                    '    bounds: [72, 72, 1136, 120]\n'
                    '    content:\n'
                    '      style: "$title"\n'
                    '      text: |\n'
                    '        <p>标题</p>\n'), {}
        monkeypatch.setattr(recipe, '_llm_chat', _llm)
        monkeypatch.setattr(
            recipe, '_run_research',
            lambda ctx: {'cards': [], 'degraded': True, 'reason': 'test'})
        import lib.slides.author as author
        monkeypatch.setattr(author, '_llm',
                            lambda messages, **kw: _llm(messages, **kw))
        # render + qa mocked: no browser, no VLM in the unit lane
        import lib.slides.recipe as r
        monkeypatch.setattr(r, '_run_render',
                            lambda ctx: {'previews': [], 'failed': [],
                                         'skipped': True})
        import lib.design_sys.visual_qa as vqa
        monkeypatch.setattr(vqa, 'visual_qa_available',
                            lambda: (False, 'test'))
        monkeypatch.setattr(
            "lib.slides._asset_preflight.prepare_deck_assets",
            lambda *a, **k: {"by_page": {}, "records": [], "findings": []})

        out = recipe.build_deck_from_topic('测试主题', str(tmp_path / 'job'))
        assert os.path.isfile(out['pptx_path'])
        assert out['pages'] == 3
        assert out['authored_pages'] == 3
        assert out['theme_id'] == 'paper-engineer'

        # Resume: a second run must skip every finished stage (checkpoint).
        n_before = calls['n']
        out2 = recipe.build_deck_from_topic('测试主题', str(tmp_path / 'job'))
        assert calls['n'] == n_before     # zero LLM calls on full resume
        assert out2['pptx_path'] == out['pptx_path']

    def test_bad_page_blocks_designer_quality_publication(
            self, tmp_path, monkeypatch):
        import lib.slides.recipe as recipe
        import lib.slides.author as author
        outline_json = ('{"title": "T", "scenario": "business-plan",'
                        ' "pages": ['
                        ' {"pageType": "cover", "key_message": "甲"},'
                        ' {"pageType": "content", "key_message": "乙"},'
                        ' {"pageType": "final", "key_message": "丙"}]}')

        def _llm(messages, **kw):
            text = messages[0]['content']
            if '大纲' in text:
                return outline_json, {}
            return '这不是 YAML: [{{{', {}
        monkeypatch.setattr(recipe, '_llm_chat', _llm)
        monkeypatch.setattr(
            recipe, '_run_research',
            lambda ctx: {'cards': [], 'degraded': True, 'reason': 'test'})
        monkeypatch.setattr(author, '_llm',
                            lambda messages, **kw: _llm(messages, **kw))
        monkeypatch.setattr(recipe, '_run_render',
                            lambda ctx: {'previews': [], 'failed': [],
                                         'skipped': True})
        import lib.design_sys.visual_qa as vqa
        monkeypatch.setattr(vqa, 'visual_qa_available',
                            lambda: (False, 'test'))
        monkeypatch.setattr(
            "lib.slides._asset_preflight.prepare_deck_assets",
            lambda *a, **k: {"by_page": {}, "records": [], "findings": []})
        from lib.production.stages import StageFailed
        with pytest.raises(StageFailed, match="stage 'author' failed"):
            recipe.build_deck_from_topic('x', str(tmp_path / 'job'))
        # The deterministic fallbacks remain as retry diagnostics, but no
        # PPTX is published under a designer-quality success status.
        assert len(list((tmp_path / 'job' / 'deck' / 'pages').glob(
            '*.page'))) == 3
        assert not list((tmp_path / 'job' / 'deck').glob('*.pptx'))

    @pytest.mark.parametrize('improves', [True, False])
    def test_layout_qa_accepts_only_measured_improvement(
            self, tmp_path, monkeypatch, improves):
        import lib.slides.recipe as recipe
        manifest = _write_deck(tmp_path, [_COVER])
        deck_dir = os.path.dirname(manifest)
        page_path = os.path.join(deck_dir, 'pages', '01.page')
        original = open(page_path, encoding='utf-8').read()
        finding = {
            'type': 'text_overflow', 'page': 1, 'elements': ['title'],
            'message': 'page 1: text "title" overflows its bounds',
        }
        initial = {'ran': True, 'ok': False, 'findings': [finding],
                   'pages': [{'index': 0, 'findings': [finding]}]}
        clean = {'ran': True, 'ok': True, 'findings': [],
                 'pages': [{'index': 0, 'findings': []}]}
        measurements = [initial, clean if improves else initial]
        if not improves:
            measurements.append(initial)  # remeasure after exact rollback

        import lib.slides.layout_qa as layout_qa
        monkeypatch.setattr(layout_qa, 'inspect_deck_layout',
                            lambda deck: measurements.pop(0))
        import lib.slides.author as author
        candidate = original.replace('[72, 200, 1136, 120]',
                                     '[72, 190, 1136, 140]')
        monkeypatch.setattr(
            author, 'author_page',
            lambda *a, **k: {'mode': 'authored', 'yaml': candidate})
        import lib.slides.render_png as render_png
        monkeypatch.setattr(render_png, 'render_page_png',
                            lambda *a, **k: str(tmp_path / 'preview.png'))
        ctx = {
            'deck_dir': deck_dir,
            'lang': 'zh',
            'model': 'kimi-k3',
            'artifacts': {
                'render': {'previews': ['dummy.png']},
                'outline': {'pages': [{'key_message': '标题'}]},
                'design': {'theme_id': 'paper-engineer'},
                'author': {},
            },
        }
        result = recipe._run_layout_qa(ctx)
        on_disk = open(page_path, encoding='utf-8').read()
        if improves:
            assert result['repaired'] == 1
            assert on_disk == candidate
        else:
            assert result['repaired'] == 0
            assert on_disk == original


# ── Runtime / starter contract ────────────────────────────

class TestRuntimeContract:
    def test_kind_matches_starter(self):
        """The kind /start dispatches on must be the kind the runtime
        registers — a mismatch means a job you can start but never poll."""
        from lib.slides.runtime import _slides_runtime
        assert _slides_runtime.kind == 'slides-deck'

    def test_start_preflight_fails_before_task_or_workdir(
            self, tmp_path, monkeypatch):
        import lib.slides.engine as engine
        import lib.slides.readiness as readiness
        import lib.slides.runtime as runtime

        monkeypatch.setattr(engine, 'slides_root', lambda: str(tmp_path))
        before = runtime._slides_runtime.task_ids()

        def unavailable():
            raise readiness.SlidesRuntimeUnavailable('playwright is absent')

        monkeypatch.setattr(readiness, 'ensure_slides_runtime_ready',
                            unavailable)
        with pytest.raises(readiness.SlidesRuntimeUnavailable):
            engine.start_slides_job('preflight-test', user_id=991)

        assert runtime._slides_runtime.task_ids() == before
        assert not list(tmp_path.iterdir())

    def test_identical_slide_starts_claim_one_task_atomically(
            self, tmp_path, monkeypatch):
        from concurrent.futures import ThreadPoolExecutor

        import lib.slides.engine as engine
        import lib.slides.readiness as readiness
        import lib.slides.runtime as runtime

        monkeypatch.setattr(engine, 'slides_root', lambda: str(tmp_path))
        monkeypatch.setattr(readiness, 'ensure_slides_runtime_ready',
                            lambda: {'cached': True})
        spawned = []
        monkeypatch.setattr(runtime._slides_runtime, 'spawn',
                            lambda task_id, *_args: spawned.append(task_id))

        def start(_index):
            return engine.start_slides_job(
                'atomic-slide-start-contract', user_id=992)

        with ThreadPoolExecutor(max_workers=8) as pool:
            results = list(pool.map(start, range(24)))

        task_ids = {result['task_id'] for result in results}
        assert len(task_ids) == 1
        assert sum(not result['deduped'] for result in results) == 1
        assert spawned == [next(iter(task_ids))]
        assert [path.name for path in tmp_path.iterdir()] == ['jobs']

    def test_stage_abort_settles_task_as_aborted_not_error(
            self, monkeypatch):
        import lib.slides.engine as engine
        import lib.slides.recipe as recipe
        import lib.slides.runtime as runtime
        from lib.production.stages import StageAborted

        states = []
        finishes = []
        monkeypatch.setattr(engine, '_write_manifest',
                            lambda _task, state: states.append(state))
        monkeypatch.setattr(engine, '_emit', lambda *_args: None)
        monkeypatch.setattr(runtime._slides_runtime, 'mark_running',
                            lambda _task_id: None)
        monkeypatch.setattr(
            runtime._slides_runtime, 'finish',
            lambda task_id, **kwargs: finishes.append((task_id, kwargs)))

        def abort_build(*_args, **_kwargs):
            raise StageAborted('cancelled')

        monkeypatch.setattr(recipe, 'build_deck_from_topic', abort_build)

        engine.run_slides_task({
            'task_id': 'slides-test', 'topic': 'x', 'workdir': '/tmp/x',
            'user_id': 1})

        assert states == ['running', 'aborted']
        assert finishes == [(
            'slides-test',
            {'error': 'aborted', 'error_context': 'slides:abort'})]

    def test_quality_gate_failure_is_actionable_and_retryable(
            self, monkeypatch):
        import lib.slides.engine as engine
        import lib.slides.readiness as readiness
        import lib.slides.recipe as recipe
        import lib.slides.runtime as runtime
        from lib.production.stages import StageFailed

        states = []
        finishes = []
        monkeypatch.setattr(engine, '_write_manifest',
                            lambda _task, state: states.append(state))
        monkeypatch.setattr(engine, '_emit', lambda *_args: None)
        monkeypatch.setattr(readiness, 'ensure_slides_runtime_ready',
                            lambda: {'cached': True})
        monkeypatch.setattr(runtime._slides_runtime, 'mark_running',
                            lambda _task_id: None)
        monkeypatch.setattr(
            runtime._slides_runtime, 'finish',
            lambda task_id, **kwargs: finishes.append((task_id, kwargs)))
        monkeypatch.setattr(
            recipe, 'build_deck_from_topic',
            lambda *_args, **_kwargs: (_ for _ in ()).throw(StageFailed(
                'author', 'gate rejected the artifact',
                ['2 of 12 pages used a fallback'])))

        engine.run_slides_task({
            'task_id': 'slides-quality', 'topic': 'x', 'workdir': '/tmp/x',
            'user_id': 1})

        assert states == ['running', 'error']
        error = finishes[0][1]['error']
        assert error['retryable'] is True
        assert error['context'] == 'slides:quality:author'
        assert '2 of 12 pages' in error['detail']

    def test_preview_renderer_reuses_browser_waits_on_assets_and_aborts(
            self, tmp_path, monkeypatch):
        import threading
        from pathlib import Path

        import lib.slides.render_png as render_png

        pytest.importorskip('playwright')
        deck = parse_deck(_write_deck(tmp_path, [_COVER, _COVER]))
        monkeypatch.setattr(render_png, 'render_page_html',
                            lambda *_args, **_kwargs: '<html>page</html>')
        monkeypatch.setattr('chromium_env.ensure_chromium_env',
                            lambda _env: None)
        abort_event = threading.Event()
        state = {'launches': 0, 'evaluates': 0, 'screenshots': 0,
                 'abort_after_first': False, 'closed': 0}

        class BrowserPage:
            def goto(self, *_args, **_kwargs):
                return None

            def evaluate(self, script, limit):
                assert 'document.fonts.ready' in script
                assert 100 <= limit <= 2000
                state['evaluates'] += 1

            def screenshot(self, *, path, type):
                assert type == 'png'
                Path(path).write_bytes(b'png')
                state['screenshots'] += 1
                if state['abort_after_first']:
                    abort_event.set()

        class Browser:
            def new_page(self, **_kwargs):
                return BrowserPage()

            def close(self):
                state['closed'] += 1

        class Chromium:
            def launch(self):
                state['launches'] += 1
                return Browser()

        class Playwright:
            chromium = Chromium()

        class Manager:
            def __enter__(self):
                return Playwright()

            def __exit__(self, *_args):
                return False

        monkeypatch.setattr('playwright.sync_api.sync_playwright', Manager)

        complete = render_png.render_previews(
            deck, str(tmp_path / 'complete'))
        assert len(complete['pages']) == 2
        assert state['launches'] == 1
        assert state['evaluates'] == 2
        assert state['screenshots'] == 2
        assert not list((tmp_path / 'complete').rglob('*.html'))

        state['abort_after_first'] = True
        abort_event.clear()
        with pytest.raises(InterruptedError):
            render_png.render_previews(
                deck, str(tmp_path / 'aborted'),
                abort_check=abort_event.is_set)
        assert state['launches'] == 2
        assert state['screenshots'] == 3
        assert (tmp_path / 'aborted/pages/01.png').is_file()
        assert not (tmp_path / 'aborted/pages/02.png').exists()
        assert state['closed'] == 2

        oversized = pptd.Deck(
            title='too large', size=(10_000, 720), theme=deck.theme,
            pages=deck.pages, root=deck.root)
        with pytest.raises(ValueError, match='finite limit'):
            render_png.render_previews(
                oversized, str(tmp_path / 'oversized'))
        assert state['launches'] == 2

        monkeypatch.setattr(render_png, '_MAX_PREVIEW_PNG_BYTES', 2)
        state['abort_after_first'] = False
        abort_event.clear()
        bounded = render_png.render_previews(
            deck, str(tmp_path / 'bounded'))
        assert bounded['pages'] == []
        assert len(bounded['failed']) == 2
        assert not list((tmp_path / 'bounded').rglob('*.png'))

    def test_slide_request_geometry_is_one_finite_contract(self, tmp_path):
        from lib.slides.contracts import (normalise_slide_briefs,
                                          normalise_slide_image_references,
                                          normalise_slide_page_count,
                                          normalise_slide_size,
                                          normalise_slide_topic)
        from lib.slides.recipe import build_deck_from_topic

        assert normalise_slide_page_count(999) == 20
        assert normalise_slide_page_count('bad') == 12
        assert normalise_slide_size('960x540') == (960, 540)
        assert normalise_slide_size([720, 540]) == (720, 540)
        with pytest.raises(ValueError, match='unsupported slide size'):
            normalise_slide_size((50_000, 50_000))
        brief = normalise_slide_briefs(
            [{'pageType': 'unknown', 'content_notes': 'x' * 5000}],
            maximum=20)[0]
        assert brief['pageType'] == 'content'
        assert len(brief['content_notes']) == 4000
        with pytest.raises(ValueError, match='4000-character'):
            normalise_slide_topic('x' * 4001)
        consumed = []

        def references():
            for index in range(1000):
                consumed.append(index)
                yield f'https://example.test/{index}.png'

        image_refs, image_findings = normalise_slide_image_references(
            references())
        assert len(image_refs) == 64
        assert len(consumed) == 65
        assert len(image_findings) == 1

        workdir = tmp_path / 'invalid-size-job'
        with pytest.raises(ValueError, match='unsupported slide size'):
            build_deck_from_topic('x', str(workdir), size=(50_000, 50_000))
        assert not workdir.exists()

    def test_task_persists_requested_model_and_owner(self, tmp_path):
        from lib.slides.runtime import _new_slides_task, _slides_task_id
        owner_user_id = 7
        task = _new_slides_task(
            _slides_task_id(), topic='小米澎程', workdir=str(tmp_path),
            lang='zh', style='brand film', max_pages=8, size=(1280, 720),
            model='kimi-k3', user_id=owner_user_id)
        assert task['model'] == 'kimi-k3'
        assert task['user_id'] == owner_user_id

    def test_local_input_images_are_copied_into_deck(self, tmp_path):
        from lib.slides.recipe import _materialise_input_images

        source = tmp_path / 'outside' / 'hero.png'
        source.parent.mkdir()
        source.write_bytes(b'png-ish test bytes')
        deck_dir = tmp_path / 'job' / 'deck'
        deck_dir.mkdir(parents=True)

        paths, findings = _materialise_input_images(
            [str(source), 'https://example.test/reference.png'], str(deck_dir))

        assert findings == []
        assert paths[0].startswith('media/input_01_')
        assert (deck_dir / paths[0]).read_bytes() == source.read_bytes()
        assert paths[1] == 'https://example.test/reference.png'

    def test_missing_local_input_is_omitted_with_finding(self, tmp_path):
        from lib.slides.recipe import _materialise_input_images

        paths, findings = _materialise_input_images(
            ['/definitely/missing/hero.png'], str(tmp_path))

        assert paths == []
        assert len(findings) == 1

    def test_oversized_local_input_is_rejected_before_reading(self, tmp_path):
        from lib.slides._media_io import MAX_SLIDE_IMAGE_BYTES
        from lib.slides.recipe import _materialise_input_images

        source = tmp_path / 'oversized.png'
        with open(source, 'wb') as output:
            output.truncate(MAX_SLIDE_IMAGE_BYTES + 1)

        paths, findings = _materialise_input_images(
            [str(source)], str(tmp_path / 'deck'))

        assert paths == []
        assert len(findings) == 1

    def test_assets_localise_background_only_page_with_bounded_stream(
            self, tmp_path, monkeypatch):
        import contextlib

        import lib.slides.recipe as recipe

        page = '''pageType: content
background:
  type: image
  src: https://example.test/background.png
elements: []
'''
        manifest = _write_deck(tmp_path, [page])
        deck = parse_deck(manifest)

        class Response:
            status_code = 200
            headers = {'Content-Length': '1024', 'Content-Type': 'image/png'}

            @staticmethod
            def iter_content(*, chunk_size):
                assert chunk_size > 0
                yield b'x' * 1024

        @contextlib.contextmanager
        def stream(*_args, **_kwargs):
            calls.append(1)
            yield Response()

        calls = []
        monkeypatch.setattr('lib.http_client.http_stream', stream)
        ctx = {
            'deck_dir': deck.root,
            'size': deck.size,
            'artifacts': {
                'outline': {'title': deck.title, 'pages': [{}]},
                'design': {'theme_tokens': deck.theme},
                'author': {'page_files': ['pages/01.page']},
            },
        }

        result = recipe._run_assets(ctx)

        assert result['downloaded'] == 1
        assert result['downloaded_bytes'] == 1024
        remote_files = list((tmp_path / 'deck' / 'media').glob('remote_*.png'))
        assert len(remote_files) == 1
        assert remote_files[0].stat().st_size == 1024
        page_path = tmp_path / 'deck' / 'pages/01.page'
        text = page_path.read_text(
            encoding='utf-8')
        assert f'media/{remote_files[0].name}' in text

        # Simulate a crash before the page rewrite was checkpointed: restoring
        # the remote source must reuse the exact cached bytes without HTTP.
        page_path.write_text(page, encoding='utf-8')
        second = recipe._run_assets(ctx)
        assert second['downloaded'] == 0
        assert second['reused'] == 1
        assert calls == [1]

        # Same-size corruption fails the SHA-256 check and redownloads once.
        remote_files[0].write_bytes(b'z' * 1024)
        page_path.write_text(page, encoding='utf-8')
        third = recipe._run_assets(ctx)
        assert third['downloaded'] == 1
        assert calls == [1, 1]

    def test_topic_builder_normalises_relative_workdir(self, tmp_path,
                                                       monkeypatch):
        import lib.slides.recipe as recipe

        monkeypatch.chdir(tmp_path)
        seen = {}

        def fake_stages(stages, ctx, **kwargs):
            seen['workdir'] = ctx['workdir']
            return {
                'outline': {'title': 'T', 'scenario': 'brand-creative',
                            'pages': []},
                'design': {'theme_id': 'editorial-ink'},
                'author': {'total': 0, 'authored': 0},
                'render': {'previews': []},
                'visual_qa': {},
                'export': {'pptx_path': str(tmp_path / 'x.pptx'), 'bytes': 0},
            }

        monkeypatch.setattr(recipe, 'run_stages', fake_stages)
        recipe.build_deck_from_topic('x', 'relative-slide-job')

        assert seen['workdir'] == str(tmp_path / 'relative-slide-job')

    def test_fallback_page_validates(self, tmp_path):
        from lib.slides.author import fallback_page
        from lib.design_sys.themes import get_theme
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        theme = get_theme('paper-engineer')
        yaml_text = fallback_page(deck, {'key_message': '完整判断句'},
                                  theme=theme)
        page_path = os.path.join(deck.root, 'pages', '99_fb.page')
        open(page_path, 'w').write(yaml_text)
        deck2 = parse_deck(deck.manifest_path)
        deck2.pages.append(pptd.Page(path='pages/99_fb.page',
                                     elements=__import__('yaml').safe_load(
                                         yaml_text)['elements'],
                                     background={'type': 'solid',
                                                 'color': '$bg'}))
        assert validate_deck(deck2) == []

    def test_internal_outline_copy_and_source_ids_are_rejected(self, tmp_path):
        from lib.slides.author import _validate_page_text
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        leaked = _COVER.replace(
            '<p><strong>一寸万象</strong></p>',
            '<p>章节引导页，引出后续时间线。[S1]</p>')

        findings = _validate_page_text(deck, 'pages/leaked.page', leaked)

        assert any('内部策划' in finding for finding in findings)
        assert any('内部来源编号' in finding for finding in findings)

    def test_chart_numbers_require_page_evidence(self):
        from lib.slides.author import _brief_fidelity_findings
        page = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: invented
    elementType: chart
    bounds: [72, 180, 800, 400]
    chartType: bar
    data:
      categories: [A, B]
      series:
        - {name: Model, values: [88, 85]}
'''
        findings = _brief_fidelity_findings(
            {'key_message': 'The model is competitive.',
             'content_notes': 'No exact benchmark score was published.'},
            page, page_index=0, total=12)

        assert findings and '88' in findings[0] and '85' in findings[0]

    def test_runtime_readiness_launches_browser_once(self, monkeypatch):
        import types

        import chromium_env
        import lib.slides.readiness as readiness

        calls = {'launch': 0, 'close': 0, 'native': 0}

        class Browser:
            version = 'test-browser'

            def close(self):
                calls['close'] += 1

        class Chromium:
            @staticmethod
            def launch(*, headless):
                assert headless is True
                calls['launch'] += 1
                return Browser()

        class Manager:
            def __enter__(self):
                return types.SimpleNamespace(chromium=Chromium())

            def __exit__(self, *_args):
                return False

        fake = types.ModuleType('playwright.sync_api')
        fake.sync_playwright = Manager
        monkeypatch.setitem(sys.modules, 'playwright.sync_api', fake)
        monkeypatch.setattr(
            chromium_env, 'ensure_chromium_env',
            lambda _env: calls.__setitem__('native', calls['native'] + 1))
        monkeypatch.setattr(readiness, '_ready_identity', None)

        first = readiness.ensure_slides_runtime_ready()
        second = readiness.ensure_slides_runtime_ready()

        assert first['cached'] is False and second['cached'] is True
        assert calls == {'launch': 1, 'close': 1, 'native': 1}

    def test_runtime_readiness_rejects_missing_export_dependency(
            self, monkeypatch):
        import lib.slides.readiness as readiness

        native_import = readiness.importlib.import_module

        def import_module(name):
            if name == 'brotli':
                raise ModuleNotFoundError("No module named 'brotli'")
            return native_import(name)

        monkeypatch.setattr(readiness, '_ready_identity', None)
        monkeypatch.setattr(readiness.importlib, 'import_module', import_module)

        with pytest.raises(readiness.SlidesRuntimeUnavailable) as caught:
            readiness.ensure_slides_runtime_ready()

        assert 'required export dependency brotli' in str(caught.value)

    def test_tool_schema_and_family(self):
        from lib.tools.produce import (PRODUCE_SLIDES_TOOL,
                                       PRODUCE_SLIDES_TOOL_NAME,
                                       PRODUCE_TOOL_NAMES)
        assert PRODUCE_SLIDES_TOOL_NAME in PRODUCE_TOOL_NAMES
        assert PRODUCE_SLIDES_TOOL['function']['name'] == 'produce_slides'
        assert 'topic' in PRODUCE_SLIDES_TOOL['function']['parameters'][
            'required']
        assert 'model' in PRODUCE_SLIDES_TOOL['function']['parameters'][
            'properties']


# ── P4: native chart / font embedding / import round-trip ──

_CHART_PAGE = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: c
    elementType: chart
    bounds: [80, 120, 600, 360]
    chartType: bar
    data:
      categories: ["Q1", "Q2", "Q3"]
      series:
        - name: "营收"
          values: [10, 20, 15]
'''


class TestP4:
    def test_missing_fonttools_degrades_without_runtime_install(
            self, tmp_path, monkeypatch):
        """A request must never invoke pip or mutate the serving environment."""
        import builtins

        import lib.design_sys.fonts as _fonts
        from lib.slides.export_pptx import _font_file_for_embedding

        font_path = tmp_path / 'font.ttf'
        font_path.write_bytes(b'fake-font')
        monkeypatch.setattr(
            _fonts, 'ensure_font', lambda _font_id, _weight: str(font_path))

        real_import = builtins.__import__

        def _without_fonttools(name, *args, **kwargs):
            if name == 'fontTools.ttLib':
                raise ImportError('injected missing dependency')
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, '__import__', _without_fonttools)
        assert _font_file_for_embedding('missing', 400) == (
            b'', False, False)

    def test_chart_is_native_ooxml(self, tmp_path):
        """A chart element must export as a real OOXML chart part (selectable
        in PowerPoint), not a flattened image."""
        pytest.importorskip('pptx')
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_CHART_PAGE]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        with zipfile.ZipFile(out) as z:
            charts = [n for n in z.namelist()
                      if n.startswith('ppt/charts/chart')]
            assert charts, 'no native chart part written'
            xml = z.read(charts[0]).decode()
        assert '<c:barChart>' in xml
        assert '营收' in xml and 'Q2' in xml

    def test_font_embedding_structure(self, tmp_path):
        """Embedded fonts: fntdata parts + rels + embeddedFontLst placed after
        notesSz (CT_Presentation order), at most one regular and one bold
        slot, and the bytes are glyf-outline TTFs (PowerPoint rejects CFF)."""
        pytest.importorskip('pptx')
        pytest.importorskip('fontTools')  # embedding degrades off without it
        import lib.design_sys.fonts as _fonts
        if not _fonts.ensure_font('misans', 400):
            pytest.skip('misans not staged locally')
        import re
        import zipfile
        from lib.slides.export_pptx import export_pptx
        deck = parse_deck(_write_deck(tmp_path, [_COVER]))
        out = str(tmp_path / 'out.pptx')
        summary = export_pptx(deck, out)
        assert summary['embeddedFonts'] >= 1
        assert summary['embeddedFontParts'] == 1
        assert summary['fontSubsetting'] is True
        assert summary['embeddedGlyphs'] < 500
        with zipfile.ZipFile(out) as z:
            assert z.testzip() is None
            parts = [n for n in z.namelist() if n.startswith('ppt/fonts/')]
            assert parts
            pres = z.read('ppt/presentation.xml').decode()
            rels = z.read('ppt/_rels/presentation.xml.rels').decode()
            ct = z.read('[Content_Types].xml').decode()
            blobs = {n: z.read(n) for n in parts}
        assert 'embedTrueTypeFonts="1"' in pres
        assert 'Extension="fntdata"' in ct
        assert '/relationships/font' in rels
        assert pres.index('<p:embeddedFontLst>') > pres.index('<p:notesSz')
        lst = re.search(r'<p:embeddedFontLst>.*?</p:embeddedFontLst>',
                        pres, re.DOTALL).group(0)
        assert '<p:bold ' in lst and '<p:regular ' not in lst
        # python-pptx must still open the re-zipped package
        from pptx import Presentation
        assert len(Presentation(out).slides) == 1
        # embedded bytes are TrueType-outline sfnt
        from fontTools.ttLib import TTFont
        import io
        for n, blob in blobs.items():
            f = TTFont(io.BytesIO(blob))
            assert 'glyf' in f, f'{n} is not a glyf-outline font'
            cmap = f.getBestCmap()
            assert all(ord(ch) in cmap for ch in '一寸万象')
            assert len(f.getGlyphOrder()) < 500

    def test_font_usage_uses_actual_family_and_style_slot(self, tmp_path):
        """An unused theme face must not embed; a one-source display face
        used as bold must occupy PowerPoint's bold slot, not regular."""
        pytest.importorskip('pptx')
        pytest.importorskip('fontTools')
        pytest.importorskip('brotli')
        import lib.design_sys.fonts as _fonts
        if not _fonts.ensure_font('smiley-sans', 400):
            pytest.skip('smiley sans not staged locally')
        page = '''pageType: content
background: {type: solid, color: "$bg"}
elements:
  - elementId: display
    elementType: text
    bounds: [80, 120, 800, 100]
    content:
      fontFamily: "Smiley Sans"
      fontSize: 48
      bold: true
      text: "澎程 A"
'''
        theme = {
            'colors': {'bg': '#FFFFFF'},
            'textStyles': {
                'unused': {'fontFamily': 'MiSans', 'fontSize': 20},
            },
        }
        deck = parse_deck(_write_deck(tmp_path, [page], theme=theme))
        from lib.slides.export_pptx import _collect_font_usage, export_pptx
        usage = _collect_font_usage(deck)
        assert set(usage) == {'Smiley Sans'}
        assert set(usage['Smiley Sans']) == {'bold'}
        assert set('澎程 A') <= usage['Smiley Sans']['bold']

        import re
        import zipfile
        out = str(tmp_path / 'actual-slot.pptx')
        summary = export_pptx(deck, out)
        assert summary['embeddedFonts'] == 1
        assert summary['embeddedFontParts'] == 1
        with zipfile.ZipFile(out) as z:
            pres = z.read('ppt/presentation.xml').decode()
        lst = re.search(r'<p:embeddedFontLst>.*?</p:embeddedFontLst>',
                        pres, re.DOTALL).group(0)
        assert 'typeface="Smiley Sans"' in lst
        assert '<p:bold ' in lst and '<p:regular ' not in lst
        assert 'MiSans' not in lst

    def test_import_round_trip_table_text(self, tmp_path):
        """Export → import: table cell text must survive the loop (the DJI
        golden deck page-17 regression — empty cells on first pass)."""
        pytest.importorskip('pptx')
        import glob
        import yaml
        from lib.slides.export_pptx import export_pptx
        from lib.slides.import_pptx import import_pptx
        deck = parse_deck(_write_deck(tmp_path, [_TABLE_PAGE]))
        out = str(tmp_path / 'out.pptx')
        export_pptx(deck, out)
        import_pptx(out, str(tmp_path / 'reimport'))
        found = []
        for p in glob.glob(str(tmp_path / 'reimport' / 'pages' / '*.page')):
            d = yaml.safe_load(open(p, encoding='utf-8'))
            for el in d.get('elements') or []:
                if el.get('elementType') == 'table':
                    for row in el.get('rows') or []:
                        found.append([
                            (c.get('text') if isinstance(c, dict) else c)
                            for c in row])
        assert ['指标', '2025'] in found
        assert ['利润', '15.8'] in found

    def test_edit_slides_tool_schema(self):
        from lib.tools.produce import (EDIT_SLIDES_TOOL,
                                       EDIT_SLIDES_TOOL_NAME,
                                       PRODUCE_TOOL_NAMES)
        assert EDIT_SLIDES_TOOL_NAME in PRODUCE_TOOL_NAMES
        fn = EDIT_SLIDES_TOOL['function']
        assert fn['name'] == 'edit_slides'
        assert set(fn['parameters']['required']) >= {'task_id', 'page',
                                                     'instruction'}
