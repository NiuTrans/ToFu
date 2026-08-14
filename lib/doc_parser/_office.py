"""lib/doc_parser/_office.py — Modern OOXML (Office 2007+) text extractors.

Provides:
  - _extract_docx  (Word 2007+, via python-docx)
  - _extract_pptx  (PowerPoint 2007+, via python-pptx)
  - _extract_xlsx  (Excel 2007+, via openpyxl)

All optional dependencies are imported lazily inside each extractor so that
importing this module never hard-fails when a backend package is missing.
"""

from lib.log import get_logger

from lib.doc_parser._truncation import truncation_warning
from lib.doc_parser._tables import render_markdown_table

logger = get_logger(__name__)

# ── .xlsx scan bounds ──
# Guard against grossly-inflated worksheet dimensions (common with embedded
# images / drawing anchors): cap the rows and columns we iterate, and bail out
# of long runs of fully-empty rows instead of walking to a phantom max_row.
_XLSX_MAX_ROWS = 1000
_XLSX_MAX_COLS = 200
_XLSX_MAX_EMPTY_RUN = 50
_XLSX_ROBUST_MAX_ROWS = 20_000
_XLSX_ROBUST_MAX_SCAN_ROWS = 50_000


def _xlsx_defined_tables(file_bytes: bytes) -> dict[str, list[tuple[str, str]]]:
    """Read Excel display names/ranges without loading a writable workbook.

    Read-only openpyxl worksheets intentionally omit ``ws.tables``. The OOXML
    relationship graph is tiny, so parse only that metadata and keep the main
    cell scan streaming. This preserves complex table names without paying the
    memory cost of ``read_only=False`` on a large workbook.
    """
    import io
    import posixpath
    import xml.etree.ElementTree as ET
    import zipfile

    rel_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    pkg_rel_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
    main_ns = 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'
    out: dict[str, list[tuple[str, str]]] = {}

    def resolve_target(base: str, target: str) -> str:
        # openpyxl currently writes package-absolute Targets (/xl/...), while
        # Excel itself commonly writes relationship-relative Targets. OOXML
        # permits both; archive member names never carry the leading slash.
        if target.startswith('/'):
            return posixpath.normpath(target).lstrip('/')
        return posixpath.normpath(posixpath.join(base, target)).lstrip('/')

    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
            workbook = ET.fromstring(archive.read('xl/workbook.xml'))
            workbook_rels = ET.fromstring(archive.read('xl/_rels/workbook.xml.rels'))
            targets = {
                rel.attrib.get('Id', ''): rel.attrib.get('Target', '')
                for rel in workbook_rels.findall(f'{{{pkg_rel_ns}}}Relationship')
            }
            for sheet in workbook.findall(f'.//{{{main_ns}}}sheet'):
                sheet_name = sheet.attrib.get('name', '')
                rel_id = sheet.attrib.get(f'{{{rel_ns}}}id', '')
                target = targets.get(rel_id, '')
                if not target:
                    continue
                sheet_path = resolve_target('xl', target)
                rel_path = posixpath.join(
                    posixpath.dirname(sheet_path), '_rels',
                    posixpath.basename(sheet_path) + '.rels')
                try:
                    sheet_xml = ET.fromstring(archive.read(sheet_path))
                    sheet_rels = ET.fromstring(archive.read(rel_path))
                except KeyError as exc:
                    logger.debug('[DocParser] missing worksheet relationship: %s',
                                 exc)
                    continue
                table_targets = {
                    rel.attrib.get('Id', ''): rel.attrib.get('Target', '')
                    for rel in sheet_rels.findall(f'{{{pkg_rel_ns}}}Relationship')
                    if rel.attrib.get('Type', '').endswith('/table')
                }
                for part in sheet_xml.findall(f'.//{{{main_ns}}}tablePart'):
                    table_target = table_targets.get(
                        part.attrib.get(f'{{{rel_ns}}}id', ''), '')
                    if not table_target:
                        continue
                    table_path = resolve_target(
                        posixpath.dirname(sheet_path), table_target)
                    try:
                        table = ET.fromstring(archive.read(table_path))
                    except KeyError as exc:
                        logger.debug('[DocParser] missing Excel table part: %s',
                                     exc)
                        continue
                    name = table.attrib.get('displayName') or table.attrib.get('name') or ''
                    ref = table.attrib.get('ref') or ''
                    if name:
                        out.setdefault(sheet_name, []).append((name, ref))
    except Exception as exc:
        logger.debug('[DocParser] Excel table metadata unavailable: %s', exc)
    return out


def _extract_docx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .docx using python-docx → Markdown-like output."""
    try:
        import docx
    except ImportError:
        logger.warning('[DocParser] python-docx not installed, cannot parse .docx')
        return {
            'text': '[python-docx not installed — run: pip install python-docx]',
            'textLength': 0,
            'totalPages': 1,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['python-docx not installed'],
        }

    import io
    warnings = []

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error('[DocParser] Failed to open .docx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .docx: {e}]',
            'textLength': 0,
            'totalPages': 1,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0
    # Denominator for any truncation warning: the full text length this doc
    # WOULD have produced. Computed from the paragraph texts already in
    # memory via python-docx, so a cut can say "kept N of M chars" rather
    # than only naming the limit it hit.
    doc_total_chars = sum(len(p.text.strip()) for p in doc.paragraphs)

    # ── Paragraphs with heading detection ──
    for para in doc.paragraphs:
        style_name = (para.style.name or '').lower() if para.style else ''
        text = para.text.strip()
        if not text:
            parts.append('')
            continue

        # Convert Word heading styles to Markdown
        if style_name.startswith('heading'):
            try:
                level = int(style_name.replace('heading', '').strip())
                level = min(max(level, 1), 6)
            except ValueError as _e_audit:
                logger.debug('[doc_parser] _extract_docx caught %s: %s', type(_e_audit).__name__, _e_audit)
                level = 2
            line = f"{'#' * level} {text}"
        elif style_name in ('title',):
            line = f"# {text}"
        elif style_name in ('subtitle',):
            line = f"## {text}"
        elif style_name.startswith('list'):
            line = f"- {text}"
        else:
            line = text

        total_chars += len(line)
        if total_chars > limit:
            remaining = limit - (total_chars - len(line))
            if remaining > 50:
                parts.append(line[:remaining])
            parts.append(f'\n[…truncated at {limit:,} chars]')
            warnings.append(truncation_warning(
                kept=total_chars - len(line) + max(remaining, 0),
                total=doc_total_chars, unit='chars',
                detail=f'char limit {limit:,}'))
            break
        parts.append(line)

    # ── Tables → Markdown tables ──
    for table in doc.tables:
        if total_chars > limit:
            break
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_md = render_markdown_table(rows)
        if table_md:
            total_chars += len(table_md)
            parts.append('')
            parts.append(table_md)

    text = '\n'.join(parts)
    # Clean up excessive blank lines
    import re
    text = re.sub(r'\n{3,}', '\n\n', text).strip()

    logger.info('[DocParser] Extracted .docx: %d paragraphs, %d tables, %s chars',
                len(doc.paragraphs), len(doc.tables), f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': max(1, len(text) // 3000),  # rough page estimate
        'isScanned': False,
        'method': 'python-docx',
        'warnings': warnings,
    }


def _extract_pptx(file_bytes: bytes, limit: int) -> dict:
    """Extract text from .pptx using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning('[DocParser] python-pptx not installed, cannot parse .pptx')
        return {
            'text': '[python-pptx not installed — run: pip install python-pptx]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['python-pptx not installed'],
        }

    import io
    warnings = []

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error('[DocParser] Failed to open .pptx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .pptx: {e}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0
    n_slides = len(prs.slides)

    for si, slide in enumerate(prs.slides, 1):
        slide_parts = [f'## Slide {si}/{n_slides}']
        for shape in slide.shapes:
            if getattr(shape, 'has_table', False):
                table = shape.table
                table_md = render_markdown_table([
                    [cell.text for cell in row.cells] for row in table.rows
                ])
                if table_md:
                    slide_parts.append(table_md)
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_parts.append(text)
        slide_text = '\n'.join(slide_parts)
        total_chars += len(slide_text)
        if total_chars > limit:
            parts.append(f'\n[…truncated at slide {si}/{n_slides}]')
            warnings.append(truncation_warning(
                kept=si - 1, total=n_slides, unit='slides',
                detail=f'stopped at slide {si}'))
            break
        parts.append(slide_text)

    text = '\n\n---\n\n'.join(parts)
    logger.info('[DocParser] Extracted .pptx: %d slides, %s chars',
                n_slides, f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': n_slides,
        'isScanned': False,
        'method': 'python-pptx',
        'warnings': warnings,
    }


def _extract_xlsx(file_bytes: bytes, limit: int, *, robust: bool = False) -> dict:
    """Extract text from .xlsx using openpyxl.

    ``robust=True`` is the persistent-knowledge path: it crosses blank gaps,
    retains multiple separated table blocks, and scans up to 20k populated
    rows / 50k physical rows per sheet. The default preserves the bounded,
    low-latency attachment behavior.
    """
    try:
        import openpyxl
    except ImportError:
        logger.warning('[DocParser] openpyxl not installed, cannot parse .xlsx')
        return {
            'text': '[openpyxl not installed — run: pip install openpyxl]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'unavailable',
            'warnings': ['openpyxl not installed'],
        }

    import io
    warnings = []

    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as e:
        logger.error('[DocParser] Failed to open .xlsx: %s', e, exc_info=True)
        return {
            'text': f'[Failed to parse .xlsx: {e}]',
            'textLength': 0,
            'totalPages': 0,
            'isScanned': False,
            'method': 'error',
            'warnings': [str(e)],
        }

    parts = []
    total_chars = 0
    n_sheets = len(wb.sheetnames)
    defined_tables = _xlsx_defined_tables(file_bytes) if robust else {}

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts = [f'## Sheet: {sheet_name}']
        if defined_tables.get(sheet_name):
            table_labels = ', '.join(
                f'{name} ({ref})' if ref else name
                for name, ref in defined_tables[sheet_name])
            sheet_parts.append(f'Defined Excel tables: {table_labels}')

        # Worksheet dimensions are often grossly inflated — embedded images,
        # drawing anchors, or stray formatting can push max_row/max_column to
        # tens of thousands even when the real data is a handful of cells.
        # Iterating the full reported range would emit millions of empty
        # cells (slow + useless output), so bound the scan up front and trim
        # trailing-empty cells / skip empty rows as we go.
        col_cap = min(ws.max_column or _XLSX_MAX_COLS, _XLSX_MAX_COLS)

        rows_data = []
        row_blocks = []
        current_block = []
        empty_run = 0
        truncated_rows = False
        # ★ Every cut must be able to report its DENOMINATOR. A warning that
        # says "truncated at 1000 rows" without saying "of 5000" gives the
        # model a numerator with no scale — it cannot tell 20% from 99%.
        rows_scanned = 0          # data rows actually walked (excl. blanks)
        empty_run_stopped_at = 0  # row index where a blank run ended the scan
        scan_capped = False
        row_cap = _XLSX_ROBUST_MAX_ROWS if robust else _XLSX_MAX_ROWS
        for physical_row, row in enumerate(
                ws.iter_rows(values_only=True, max_col=col_cap), 1):
            if robust and physical_row > _XLSX_ROBUST_MAX_SCAN_ROWS:
                scan_capped = True
                break
            cells = list(row)
            while cells and cells[-1] is None:
                cells.pop()
            if not cells:
                empty_run += 1
                if robust:
                    if current_block:
                        row_blocks.append(current_block)
                        current_block = []
                    continue
                if empty_run > _XLSX_MAX_EMPTY_RUN:
                    # This break used to be entirely SILENT. A sheet shaped
                    # "summary block / 60 blank rows / detail block" lost the
                    # whole detail block with no trace in the output at all —
                    # worse than the row cap, which at least admits it fired.
                    empty_run_stopped_at = rows_scanned
                    break
                continue
            empty_run = 0
            rows_scanned += 1
            rows_data.append(cells)
            if robust:
                current_block.append(cells)
            if len(rows_data) >= row_cap:
                truncated_rows = True
                break

        if robust and current_block:
            row_blocks.append(current_block)

        # Sheet dimensions as reported by the workbook — the denominator the
        # caller needs. Guarded because max_row/max_column can be None.
        sheet_rows = ws.max_row or 0
        sheet_cols = ws.max_column or 0

        if truncated_rows:
            warnings.append(truncation_warning(
                kept=len(rows_data), total=sheet_rows, unit='rows',
                scope=f'Sheet "{sheet_name}"',
                detail=f'row cap {row_cap:,}'))
        if scan_capped:
            warnings.append(truncation_warning(
                kept=_XLSX_ROBUST_MAX_SCAN_ROWS, total=sheet_rows, unit='rows',
                scope=f'Sheet "{sheet_name}"',
                detail='physical scan safety cap'))
        if empty_run_stopped_at and not robust:
            warnings.append(truncation_warning(
                kept=empty_run_stopped_at, total=sheet_rows, unit='rows',
                scope=f'Sheet "{sheet_name}"',
                detail=(f'stopped after {_XLSX_MAX_EMPTY_RUN} consecutive '
                        f'blank rows — content below a long blank gap is '
                        f'missing')))
        if sheet_cols > _XLSX_MAX_COLS:
            warnings.append(truncation_warning(
                kept=_XLSX_MAX_COLS, total=sheet_cols, unit='columns',
                scope=f'Sheet "{sheet_name}"'))

        if robust and row_blocks:
            for block_index, block in enumerate(row_blocks, 1):
                if len(row_blocks) > 1:
                    sheet_parts.append(f'### Table block {block_index}')
                table_md = render_markdown_table(block)
                if table_md:
                    sheet_parts.append(table_md)
        elif rows_data:
            table_md = render_markdown_table(rows_data)
            if table_md:
                sheet_parts.append(table_md)

        sheet_text = '\n'.join(sheet_parts)
        total_chars += len(sheet_text)
        if total_chars > limit:
            parts.append('\n[…truncated]')
            break
        parts.append(sheet_text)

    wb.close()
    text = '\n\n---\n\n'.join(parts)
    logger.info('[DocParser] Extracted .xlsx: %d sheets, %s chars',
                n_sheets, f'{len(text):,}')

    return {
        'text': text,
        'textLength': len(text),
        'totalPages': len(wb.sheetnames),
        'isScanned': False,
        'method': 'openpyxl',
        'warnings': warnings,
    }
