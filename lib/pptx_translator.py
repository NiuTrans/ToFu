"""lib/pptx_translator.py — PowerPoint translation with formatting preservation.

Translates all text in a .pptx file while faithfully preserving fonts, colors,
alignment, spacing, bold/italic, and table cell formatting.

Core approach:
  1. Walk every slide shape (text boxes, tables, grouped shapes).
  2. For each shape, extract per-run formatting properties.
  3. Translate the text using our existing translation infrastructure.
  4. Rebuild the shape with translated text + original formatting.

Adapted from tristan-mcinnis/PPT-Translator-Formatting-Intact-with-LLMs
(https://github.com/tristan-mcinnis/PPT-Translator-Formatting-Intact-with-LLMs).
"""
from __future__ import annotations

import time
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

from lib.log import get_logger

logger = get_logger(__name__)

__all__ = ['translate_pptx']

# ── Lazy import guard for python-pptx ──
_PPTX_AVAILABLE = None


def _check_pptx():
    """Check if python-pptx is available (cached)."""
    global _PPTX_AVAILABLE
    if _PPTX_AVAILABLE is None:
        try:
            import pptx  # noqa: F401
            _PPTX_AVAILABLE = True
        except ImportError:
            _PPTX_AVAILABLE = False
    return _PPTX_AVAILABLE


# ═══════════════════════════════════════════════════════
#  Formatting Extraction & Application
# ═══════════════════════════════════════════════════════

def _extract_run_props(run) -> dict:
    """Extract formatting properties from a single text run."""
    props = {
        'bold': None,
        'italic': None,
        'underline': None,
        'font_size': None,
        'font_name': None,
        'font_color': None,
    }
    font = run.font
    if hasattr(font, 'bold'):
        props['bold'] = font.bold
    if hasattr(font, 'italic'):
        props['italic'] = font.italic
    if hasattr(font, 'underline'):
        props['underline'] = font.underline
    if getattr(font, 'size', None) is not None:
        props['font_size'] = font.size
    if getattr(font, 'name', None):
        props['font_name'] = font.name
    if (getattr(font, 'color', None) is not None
            and getattr(font.color, 'rgb', None) is not None):
        props['font_color'] = str(font.color.rgb)
    return props


def _apply_run_props(run, props: dict):
    """Apply saved formatting properties to a text run."""
    from pptx.dml.color import RGBColor

    font = run.font
    if props.get('bold') is not None:
        font.bold = props['bold']
    if props.get('italic') is not None:
        font.italic = props['italic']
    if props.get('underline') is not None:
        font.underline = props['underline']
    if props.get('font_size') is not None:
        font.size = props['font_size']  # already in EMU
    if props.get('font_name'):
        font.name = props['font_name']
    if props.get('font_color'):
        try:
            font.color.rgb = RGBColor.from_string(props['font_color'])
        except Exception as e:
            logger.debug('[PPTXTranslator] Could not set font color %s: %s',
                         props['font_color'], e)


def _extract_paragraph_props(paragraph) -> dict:
    """Extract paragraph-level formatting."""
    props = {
        'alignment': None,
        'line_spacing': None,
        'space_before': None,
        'space_after': None,
        'level': getattr(paragraph, 'level', 0),
    }
    if getattr(paragraph, 'alignment', None) is not None:
        props['alignment'] = paragraph.alignment
    if getattr(paragraph, 'line_spacing', None) is not None:
        props['line_spacing'] = paragraph.line_spacing
    if getattr(paragraph, 'space_before', None) is not None:
        props['space_before'] = paragraph.space_before
    if getattr(paragraph, 'space_after', None) is not None:
        props['space_after'] = paragraph.space_after
    return props


def _apply_paragraph_props(paragraph, props: dict):
    """Apply saved paragraph-level formatting."""
    if props.get('alignment') is not None:
        paragraph.alignment = props['alignment']
    if props.get('line_spacing') is not None:
        try:
            paragraph.line_spacing = props['line_spacing']
        except Exception as e:
            logger.debug('[PPTXTranslator] Could not set line_spacing: %s', e)
    if props.get('space_before') is not None:
        try:
            paragraph.space_before = props['space_before']
        except Exception as e:
            logger.debug('[PPTXTranslator] Could not set space_before: %s', e)
    if props.get('space_after') is not None:
        try:
            paragraph.space_after = props['space_after']
        except Exception as e:
            logger.debug('[PPTXTranslator] Could not set space_after: %s', e)
    if props.get('level') is not None:
        try:
            paragraph.level = props['level']
        except Exception as e:
            logger.debug('[PPTXTranslator] Could not set level: %s', e)


# ═══════════════════════════════════════════════════════
#  Shape-level text extraction & translation
# ═══════════════════════════════════════════════════════

def _extract_text_frame_data(text_frame) -> List[dict]:
    """Extract all paragraphs with their runs and formatting from a text frame.

    Returns a list of paragraph dicts, each containing:
        - 'para_props': paragraph-level formatting
        - 'runs': list of {'text': str, 'props': dict}
        - 'full_text': concatenated paragraph text
    """
    paragraphs = []
    for para in text_frame.paragraphs:
        para_data = {
            'para_props': _extract_paragraph_props(para),
            'runs': [],
            'full_text': '',
        }
        for run in para.runs:
            para_data['runs'].append({
                'text': run.text,
                'props': _extract_run_props(run),
            })
        para_data['full_text'] = ''.join(r['text'] for r in para_data['runs'])
        paragraphs.append(para_data)
    return paragraphs


def _rebuild_text_frame(text_frame, translated_paragraphs: List[dict]):
    """Rebuild a text frame with translated text, preserving per-run formatting.

    Strategy:
      - If original had multiple runs (mixed formatting), we try to
        distribute translated text proportionally across runs.
      - If original had a single run, put all translated text there.
      - Paragraph-level formatting (alignment, spacing) is preserved.
    """
    from pptx.oxml.ns import qn

    # Clear existing paragraphs (keep the first one, remove extras)
    # python-pptx requires at least one paragraph
    while len(text_frame.paragraphs) > 1:
        p_elem = text_frame.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    for i, para_data in enumerate(translated_paragraphs):
        if i == 0:
            para = text_frame.paragraphs[0]
            # Clear existing runs
            for run in list(para.runs):
                run._r.getparent().remove(run._r)
            # Also remove any direct text nodes (a:r elements)
            p_elem = para._p
            for child in list(p_elem):
                if child.tag.endswith('}r'):
                    p_elem.remove(child)
        else:
            # Add new paragraph (SubElement attaches it to text_frame._txBody)
            from lxml import etree
            etree.SubElement(text_frame._txBody, qn('a:p'))
            para = text_frame.paragraphs[-1]

        # Apply paragraph-level props
        _apply_paragraph_props(para, para_data['para_props'])

        translated_text = para_data.get('translated_text', para_data['full_text'])
        original_runs = para_data['runs']

        if not translated_text.strip() and not any(r['text'].strip() for r in original_runs):
            # Empty paragraph — keep it empty
            continue

        if len(original_runs) <= 1:
            # Single run (or no runs) — simple case
            run = para.add_run()
            run.text = translated_text
            if original_runs:
                _apply_run_props(run, original_runs[0]['props'])
        else:
            # Multiple runs with different formatting.
            # Distribute translated text proportionally.
            _distribute_text_across_runs(para, translated_text, original_runs)


def _distribute_text_across_runs(paragraph, translated_text: str, original_runs: List[dict]):
    """Distribute translated text proportionally across multiple formatted runs.

    When the original had runs like:
        [bold:"Hello "][normal:"world"][italic:"!"]
    We distribute the translated text proportionally to maintain formatting
    boundaries approximately.
    """
    total_orig_len = sum(len(r['text']) for r in original_runs)
    if total_orig_len == 0:
        # All empty runs — just put text in first run
        run = paragraph.add_run()
        run.text = translated_text
        if original_runs:
            _apply_run_props(run, original_runs[0]['props'])
        return

    # Calculate proportional lengths
    remaining = translated_text
    for idx, orig_run in enumerate(original_runs):
        run = paragraph.add_run()
        if idx == len(original_runs) - 1:
            # Last run gets all remaining text
            run.text = remaining
        else:
            proportion = len(orig_run['text']) / total_orig_len
            chars = max(1, round(len(translated_text) * proportion))
            # Try to break at a space boundary
            cut_point = min(chars, len(remaining))
            # Look for a space near the cut point
            if cut_point < len(remaining):
                space_pos = remaining.rfind(' ', max(0, cut_point - 10), cut_point + 10)
                if space_pos > 0:
                    cut_point = space_pos + 1
            run.text = remaining[:cut_point]
            remaining = remaining[cut_point:]
        _apply_run_props(run, orig_run['props'])


# ═══════════════════════════════════════════════════════
#  Table handling
# ═══════════════════════════════════════════════════════

def _extract_table_data(table) -> List[List[dict]]:
    """Extract table cell text + formatting."""
    rows = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            cell_data = {
                'paragraphs': _extract_text_frame_data(cell.text_frame),
                'margin_left': cell.margin_left,
                'margin_right': cell.margin_right,
                'margin_top': cell.margin_top,
                'margin_bottom': cell.margin_bottom,
                'vertical_anchor': cell.vertical_anchor,
            }
            row_data.append(cell_data)
        rows.append(row_data)
    return rows


def _apply_translated_table(table, translated_rows: List[List[dict]]):
    """Apply translated text back to table cells with formatting."""
    for row_idx, row in enumerate(table.rows):
        if row_idx >= len(translated_rows):
            break
        for col_idx, cell in enumerate(row.cells):
            if col_idx >= len(translated_rows[row_idx]):
                break
            cell_data = translated_rows[row_idx][col_idx]
            try:
                # Restore margins
                if cell_data.get('margin_left') is not None:
                    cell.margin_left = cell_data['margin_left']
                if cell_data.get('margin_right') is not None:
                    cell.margin_right = cell_data['margin_right']
                if cell_data.get('margin_top') is not None:
                    cell.margin_top = cell_data['margin_top']
                if cell_data.get('margin_bottom') is not None:
                    cell.margin_bottom = cell_data['margin_bottom']
                if cell_data.get('vertical_anchor') is not None:
                    cell.vertical_anchor = cell_data['vertical_anchor']
                # Rebuild text
                _rebuild_text_frame(cell.text_frame, cell_data['paragraphs'])
            except Exception as e:
                logger.warning('[PPTXTranslator] Error applying cell [%d][%d]: %s',
                               row_idx, col_idx, e)


# ═══════════════════════════════════════════════════════
#  Group shape recursion
# ═══════════════════════════════════════════════════════

def _process_grouped_shapes(group_shape, translate_fn: Callable[[str], str],
                            stats: dict):
    """Recursively process shapes inside a group."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in group_shape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _process_grouped_shapes(shape, translate_fn, stats)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            _translate_table_shape(shape.table, translate_fn, stats)
        elif hasattr(shape, 'text_frame'):
            _translate_text_shape(shape, translate_fn, stats)


# ═══════════════════════════════════════════════════════
#  Per-shape translation
# ═══════════════════════════════════════════════════════

def _translate_text_shape(shape, translate_fn: Callable[[str], str],
                          stats: dict):
    """Translate a text shape in-place, preserving formatting."""
    if not hasattr(shape, 'text_frame'):
        return
    if not shape.text_frame.text.strip():
        return

    paragraphs = _extract_text_frame_data(shape.text_frame)

    for para_data in paragraphs:
        original_text = para_data['full_text']
        if not original_text.strip():
            para_data['translated_text'] = original_text
            continue
        try:
            translated = translate_fn(original_text)
            para_data['translated_text'] = translated
            stats['chars_translated'] += len(original_text)
            stats['segments'] += 1
        except Exception as e:
            logger.warning('[PPTXTranslator] Shape translation failed (keeping original): %s',
                           str(e)[:200])
            para_data['translated_text'] = original_text
            stats['errors'] += 1

    try:
        _rebuild_text_frame(shape.text_frame, paragraphs)
    except Exception as e:
        logger.error('[PPTXTranslator] Failed to rebuild text frame: %s', e, exc_info=True)
        stats['errors'] += 1


def _translate_table_shape(table, translate_fn: Callable[[str], str],
                           stats: dict):
    """Translate all cells in a table, preserving formatting."""
    table_data = _extract_table_data(table)

    for row in table_data:
        for cell_data in row:
            for para_data in cell_data['paragraphs']:
                original_text = para_data['full_text']
                if not original_text.strip():
                    para_data['translated_text'] = original_text
                    continue
                try:
                    translated = translate_fn(original_text)
                    para_data['translated_text'] = translated
                    stats['chars_translated'] += len(original_text)
                    stats['segments'] += 1
                except Exception as e:
                    logger.warning('[PPTXTranslator] Table cell translation failed: %s',
                                   str(e)[:200])
                    para_data['translated_text'] = original_text
                    stats['errors'] += 1

    try:
        _apply_translated_table(table, table_data)
    except Exception as e:
        logger.error('[PPTXTranslator] Failed to rebuild table: %s', e, exc_info=True)
        stats['errors'] += 1


# ═══════════════════════════════════════════════════════
#  Translation cache
# ═══════════════════════════════════════════════════════

class _TranslationCache:
    """Simple thread-safe cache to avoid re-translating identical strings."""

    def __init__(self, translate_fn: Callable[[str], str]):
        self._fn = translate_fn
        self._cache: Dict[str, str] = {}

    def translate(self, text: str) -> str:
        if not text or text.isspace():
            return text
        if text in self._cache:
            return self._cache[text]
        result = self._fn(text)
        self._cache[text] = result
        return result


# ═══════════════════════════════════════════════════════
#  Main entry point
# ═══════════════════════════════════════════════════════

def translate_pptx(
    input_path: str | Path,
    output_path: str | Path | None = None,
    *,
    translate_fn: Callable[[str], str],
    progress_fn: Optional[Callable[[int, int, str], None]] = None,
) -> dict:
    """Translate a .pptx file, preserving all formatting.

    Args:
        input_path: Path to the source .pptx file.
        output_path: Path for the translated .pptx. If None, auto-generated
            as ``<stem>_translated.pptx`` next to the input.
        translate_fn: Callable that translates a text string. Receives a single
            string argument and returns the translated string. This is typically
            wired to our existing ``_translate_one_chunk()`` or MT provider.
        progress_fn: Optional callback ``(current_slide, total_slides, status_msg)``.

    Returns:
        Dict with keys:
            ok: bool
            output_path: str  (path to translated file)
            slides: int
            segments: int  (number of text segments translated)
            chars_translated: int
            errors: int
            elapsed: float
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if not _check_pptx():
        return {
            'ok': False,
            'error': 'python-pptx not installed — run: pip install python-pptx',
        }

    from pptx import Presentation

    input_path = Path(input_path)
    if not input_path.is_file():
        return {'ok': False, 'error': f'File not found: {input_path}'}
    if input_path.suffix.lower() not in ('.pptx',):
        return {'ok': False, 'error': f'Unsupported format: {input_path.suffix} (only .pptx supported)'}

    if output_path is None:
        output_path = input_path.parent / f'{input_path.stem}_translated.pptx'
    output_path = Path(output_path)

    t0 = time.time()
    stats = {'segments': 0, 'chars_translated': 0, 'errors': 0}

    # Wrap translate_fn with cache to avoid re-translating repeated text
    cached_translate = _TranslationCache(translate_fn)

    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        logger.error('[PPTXTranslator] Failed to open %s: %s', input_path.name, e, exc_info=True)
        return {'ok': False, 'error': f'Failed to open PPTX: {e}'}

    n_slides = len(prs.slides)
    logger.info('[PPTXTranslator] Processing %s: %d slides', input_path.name, n_slides)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        if progress_fn:
            try:
                progress_fn(slide_idx, n_slides, f'Translating slide {slide_idx}/{n_slides}')
            except Exception:
                pass

        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    _process_grouped_shapes(shape, cached_translate.translate, stats)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    _translate_table_shape(shape.table, cached_translate.translate, stats)
                elif hasattr(shape, 'text_frame'):
                    _translate_text_shape(shape, cached_translate.translate, stats)
            except Exception as e:
                logger.warning('[PPTXTranslator] Shape error on slide %d: %s',
                               slide_idx, e)
                stats['errors'] += 1

        logger.debug('[PPTXTranslator] Slide %d/%d done — segments=%d chars=%d',
                     slide_idx, n_slides, stats['segments'], stats['chars_translated'])

    # Save
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
    except Exception as e:
        logger.error('[PPTXTranslator] Failed to save %s: %s', output_path, e, exc_info=True)
        return {'ok': False, 'error': f'Failed to save translated file: {e}'}

    elapsed = time.time() - t0
    logger.info('[PPTXTranslator] Done: %s → %s — %d slides, %d segments, %d chars, %d errors, %.1fs',
                input_path.name, output_path.name, n_slides, stats['segments'],
                stats['chars_translated'], stats['errors'], elapsed)

    return {
        'ok': True,
        'output_path': str(output_path),
        'slides': n_slides,
        'segments': stats['segments'],
        'chars_translated': stats['chars_translated'],
        'errors': stats['errors'],
        'elapsed': round(elapsed, 1),
    }


def translate_pptx_bytes(
    file_bytes: bytes,
    filename: str = 'presentation.pptx',
    *,
    translate_fn: Callable[[str], str],
    progress_fn: Optional[Callable[[int, int, str], None]] = None,
) -> Tuple[Optional[bytes], dict]:
    """Translate a PPTX from bytes, returning translated bytes.

    Convenience wrapper that writes to a temp file, translates,
    reads back the result bytes.

    Args:
        file_bytes: Raw PPTX file content.
        filename: Original filename (used for logging and output naming).
        translate_fn: Translation function.
        progress_fn: Optional progress callback.

    Returns:
        (translated_bytes_or_None, result_dict)
    """
    import tempfile

    if not _check_pptx():
        return None, {
            'ok': False,
            'error': 'python-pptx not installed — run: pip install python-pptx',
        }

    stem = Path(filename).stem
    with tempfile.TemporaryDirectory(prefix='pptx_translate_') as tmpdir:
        input_path = Path(tmpdir) / filename
        output_path = Path(tmpdir) / f'{stem}_translated.pptx'

        input_path.write_bytes(file_bytes)
        result = translate_pptx(
            input_path,
            output_path,
            translate_fn=translate_fn,
            progress_fn=progress_fn,
        )

        if result.get('ok') and output_path.is_file():
            return output_path.read_bytes(), result
        return None, result
